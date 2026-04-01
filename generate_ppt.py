from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_FILE = "Tikur_Abay_Logistics_System.pptx"
LOGO_PATH = Path("/Users/getnetbelay/Documents/Tikur_Abay/apps/admin/public/branding/tikur-abay-logo.png")

PHASE_COLORS = {
    "Commercial": RGBColor(20, 84, 120),
    "Origin": RGBColor(11, 114, 105),
    "Shipping": RGBColor(35, 95, 168),
    "Control": RGBColor(180, 106, 28),
    "Finance": RGBColor(124, 72, 163),
    "Delivery": RGBColor(40, 118, 66),
    "Closure": RGBColor(69, 84, 100),
}

SLIDES = [
    {
        "title": "Tikur Abay End-to-End Logistics System",
        "subtitle": "Integrated shipping, finance, clearance, dispatch, and closure control",
        "phase": "Commercial",
        "cards": [
            ("Quote to closure", "One operating flow"),
            ("Cross-functional", "Shipping, bank, customs, finance"),
            ("Execution control", "Dispatch, delivery, empty return"),
        ],
    },
    {
        "title": "Executive Snapshot",
        "subtitle": "The system controls revenue, compliance, execution, and closure in one chain",
        "phase": "Commercial",
        "snapshot": [
            ("Commercial control", "Quote, booking, and carrier confirmation move from request to executable shipment."),
            ("Compliance control", "BL, manifest, LC review, and customs gates prevent uncontrolled release."),
            ("Financial control", "Charges, payment verification, and receipt issuance block release until settled."),
            ("Operational control", "Dispatch, delivery, empty return, and closure complete the shipment lifecycle."),
        ],
    },
    {
        "title": "Overview",
        "subtitle": "Full logistics lifecycle across external and internal control points",
        "phase": "Commercial",
        "diagram": [
            "Quote",
            "Booking",
            "Shipping",
            "Finance",
            "Delivery",
            "Closure",
        ],
        "notes": [
            "Integrated with shipping lines",
            "Integrated with banks (LC)",
            "Integrated with customs",
            "Integrated with inland transport",
        ],
    },
    {
        "title": "Quote Phase",
        "subtitle": "Commercial pricing starts the shipment file",
        "phase": "Commercial",
        "cards": [
            ("Trigger", "Customer requests quote"),
            ("Pricing stack", "Sea (USD)\nClearance (USD)\nInland (ETB)"),
        ],
    },
    {
        "title": "Booking & Carrier",
        "subtitle": "Carrier confirmation locks the shipping plan",
        "phase": "Commercial",
        "diagram": [
            "Booking created",
            "Carrier confirmed",
            "Vessel assigned",
            "Reference locked",
        ],
        "notes": [
            "Commercial approval becomes an executable shipment file",
        ],
    },
    {
        "title": "Origin Operations",
        "subtitle": "Physical export readiness is confirmed at source",
        "phase": "Origin",
        "diagram": [
            "Cargo ready",
            "Stuffing",
            "Container no.",
            "Seal locked",
        ],
        "notes": [
            "Origin control converts cargo readiness into a verifiable export unit",
        ],
    },
    {
        "title": "Shipping Docs",
        "subtitle": "Carrier instruction is prepared and accepted",
        "phase": "Shipping",
        "diagram": [
            "SI prepared",
            "Carrier submission",
            "Shipment confirmed",
        ],
        "notes": [
            "Document accuracy protects downstream BL, manifest, and bank processing",
        ],
    },
    {
        "title": "Shipment Execution",
        "subtitle": "Ocean movement begins after port handoff",
        "phase": "Shipping",
        "diagram": [
            "Gate-in",
            "Loaded",
            "Departure",
            "Tracking live",
        ],
        "notes": [
            "Execution status becomes visible across operations",
        ],
    },
    {
        "title": "Documentation",
        "subtitle": "Core shipping documents establish legal movement",
        "phase": "Shipping",
        "cards": [
            ("BL pack", "Master BL\nHouse BL"),
            ("Manifest", "Cargo manifest"),
        ],
    },
    {
        "title": "Bank (LC Process)",
        "subtitle": "Document pack is validated through bank review",
        "phase": "Control",
        "diagram": [
            "Docs sent",
            "Bank review",
            "Approved / Fix",
        ],
        "notes": [
            "Discrepancies loop back into document correction before release can continue",
        ],
    },
    {
        "title": "Arrival & Djibouti",
        "subtitle": "Port arrival shifts the file into release preparation",
        "phase": "Control",
        "cards": [
            ("Arrival", "Vessel arrival"),
            ("Next step", "Release preparation"),
        ],
    },
    {
        "title": "Clearance",
        "subtitle": "Customs gate controls legal release",
        "phase": "Control",
        "cards": [
            ("Process", "Customs clearance"),
            ("Outcome", "Shipment approved"),
        ],
    },
    {
        "title": "Finance",
        "subtitle": "Financial settlement controls downstream cargo release",
        "phase": "Finance",
        "diagram": [
            "Charges issued",
            "Customer pays",
            "Finance verifies",
            "Receipt issued",
        ],
        "notes": [
            "Release does not move until settlement is complete",
        ],
    },
    {
        "title": "Release Control",
        "subtitle": "Release authority is pushed into port and dry-port operations",
        "phase": "Control",
        "diagram": [
            "Receipt issued",
            "Release sent",
            "Dry port ready",
        ],
        "notes": [
            "Commercial settlement is translated into an operational release signal",
        ],
    },
    {
        "title": "Dispatch & Delivery",
        "subtitle": "Inland execution completes customer fulfillment",
        "phase": "Delivery",
        "diagram": [
            "Truck assigned",
            "Cargo released",
            "Inland move",
            "Delivered",
        ],
        "notes": [
            "Dry port to customer handoff is managed as one flow",
        ],
    },
    {
        "title": "Container Lifecycle",
        "subtitle": "Equipment return is required to complete operational control",
        "phase": "Closure",
        "diagram": [
            "Empty return",
            "Yard confirmed",
            "Container closed",
        ],
        "notes": [
            "No equipment return means the shipment is not operationally complete",
        ],
    },
    {
        "title": "Driver Process",
        "subtitle": "Driver reimbursement closes the transport expense loop",
        "phase": "Finance",
        "cards": [
            ("Field action", "Driver uploads expenses"),
            ("Finance action", "Finance reimburses"),
        ],
    },
    {
        "title": "Closure",
        "subtitle": "Shipment is only complete when every control point is satisfied",
        "phase": "Closure",
        "cards": [
            ("Status", "Shipment closed"),
            ("Result", "All processes completed"),
        ],
    },
    {
        "title": "Control Rules",
        "subtitle": "System rules enforce operating discipline",
        "phase": "Control",
        "cards": [
            ("Payment gate", "No payment -> No release"),
            ("Document gate", "No BL -> No bank process"),
            ("Customs gate", "No clearance -> No dispatch"),
            ("Closure gate", "No return -> Not closed"),
        ],
    },
]

MASTER_TIMELINE = [
    ("Commercial", "Quote\nBooking"),
    ("Origin", "Cargo ready\nStuffing"),
    ("Shipping", "SI\nVessel\nBL"),
    ("Control", "Bank\nArrival\nClearance"),
    ("Finance", "Charges\nPayment\nReceipt"),
    ("Delivery", "Release\nDispatch\nDelivery"),
    ("Closure", "Return\nClose"),
]

FLOW_STEPS = [
    "Customer requests quote",
    "Tikur Abay prepares quote",
    "Customer approves quote",
    "Booking created",
    "Carrier booking confirmed",
    "Origin preparation completed",
    "Shipping instruction prepared",
    "SI submitted to carrier",
    "Carrier confirms shipment",
    "Container loaded and vessel departs",
    "BL generated",
    "Cargo manifest generated",
    "Documents submitted to bank",
    "LC validation completed",
    "Vessel arrives Djibouti",
    "Release preparation starts",
    "Customs clearance processed",
    "Customer charges issued",
    "Customer pays",
    "Finance verifies payment",
    "Official receipt issued",
    "Release authorization sent",
    "Cargo released",
    "Truck assigned",
    "Inland transport executed",
    "Delivery completed",
    "Empty container returned",
    "Driver uploads expenses",
    "Finance reimburses driver",
    "Container closed",
    "Shipment closed",
]

FLOW_PHASES = [
    ("Commercial", FLOW_STEPS[0:5]),
    ("Origin", FLOW_STEPS[5:7]),
    ("Shipping", FLOW_STEPS[7:12]),
    ("Control", FLOW_STEPS[12:17]),
    ("Finance", FLOW_STEPS[17:21]),
    ("Delivery", FLOW_STEPS[21:26]),
    ("Closure", FLOW_STEPS[26:31]),
]


def add_logo(slide, left, top, height):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left, top, height=height)


def add_text(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_full_background(slide, phase):
    phase_color = PHASE_COLORS[phase]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(242, 245, 248)

    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.72))
    header.fill.solid()
    header.fill.fore_color.rgb = phase_color
    header.line.fill.background()

    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.42), Inches(0.95), Inches(12.45), Inches(5.95))
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(255, 255, 255)
    hero.line.color.rgb = RGBColor(216, 224, 232)

    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.56), Inches(1.1), Inches(0.14), Inches(0.52))
    strip.fill.solid()
    strip.fill.fore_color.rgb = phase_color
    strip.line.fill.background()

    add_logo(slide, Inches(11.88), Inches(0.1), Inches(0.5))


def add_title(slide, title, subtitle):
    add_text(slide, Inches(0.78), Inches(0.13), Inches(9.8), Inches(0.4), title, 24, RGBColor(255, 255, 255), bold=True)
    add_text(slide, Inches(0.82), Inches(1.08), Inches(10.9), Inches(0.45), subtitle, 19, RGBColor(24, 39, 56), bold=True)


def add_phase_chip(slide, phase):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.35), Inches(1.08), Inches(1.8), Inches(0.4))
    chip.fill.solid()
    chip.fill.fore_color.rgb = PHASE_COLORS[phase]
    chip.line.fill.background()
    p = chip.text_frame.paragraphs[0]
    p.text = phase.upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_title_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    phase_color = PHASE_COLORS[slide_data["phase"]]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(12, 24, 38)

    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(7.8), Inches(-1.6), Inches(6.4), Inches(6.4))
    glow.fill.solid()
    glow.fill.fore_color.rgb = phase_color
    glow.fill.transparency = 0.7
    glow.line.fill.background()

    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.62), Inches(12.2), Inches(5.95))
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(18, 34, 52)
    hero.line.color.rgb = RGBColor(51, 80, 109)
    hero.line.width = Pt(1.5)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(0.92), Inches(0.16), Inches(4.65))
    accent.fill.solid()
    accent.fill.fore_color.rgb = phase_color
    accent.line.fill.background()

    add_logo(slide, Inches(11.2), Inches(0.88), Inches(0.9))

    eyebrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.18), Inches(1.0), Inches(2.7), Inches(0.42))
    eyebrow.fill.solid()
    eyebrow.fill.fore_color.rgb = phase_color
    eyebrow.line.fill.background()
    eyebrow.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = eyebrow.text_frame.paragraphs[0]
    p.text = "LOGISTICS OPERATING MODEL"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    add_text(
        slide,
        Inches(1.18),
        Inches(1.72),
        Inches(6.9),
        Inches(1.55),
        "Tikur Abay\nEnd-to-End Logistics System",
        26,
        RGBColor(255, 255, 255),
        bold=True,
    )
    add_text(
        slide,
        Inches(1.18),
        Inches(3.3),
        Inches(6.3),
        Inches(0.7),
        "Integrated control across shipping, bank process, customs, finance, dispatch, and container closure.",
        17,
        RGBColor(205, 218, 229),
    )

    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.15), Inches(1.68), Inches(3.85), Inches(3.9))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = RGBColor(245, 248, 250)
    ribbon.line.fill.background()

    add_text(slide, Inches(8.52), Inches(1.95), Inches(2.7), Inches(0.3), "CORE PHASES", 12, phase_color, bold=True)
    phases = ["Quote", "Booking", "Shipping", "Finance", "Delivery", "Closure"]
    for index, label in enumerate(phases):
        top = 2.34 + index * 0.5
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.52), Inches(top), Inches(0.18), Inches(0.18))
        marker.fill.solid()
        marker.fill.fore_color.rgb = phase_color if index < 3 else RGBColor(101, 118, 135)
        marker.line.fill.background()
        add_text(slide, Inches(8.82), Inches(top - 0.06), Inches(2.2), Inches(0.28), label, 16, RGBColor(36, 54, 72), bold=True)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.72), Inches(11.6), Inches(0.5))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(29, 48, 68)
    band.line.fill.background()
    band.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = band.text_frame.paragraphs[0]
    p.text = "One operating chain from customer quote to empty-container return and shipment closure"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def add_flow_strip(slide):
    labels = ["Quote", "Booking", "Origin", "Shipping", "Control", "Finance", "Delivery", "Closure"]
    colors = [
        PHASE_COLORS["Commercial"],
        PHASE_COLORS["Commercial"],
        PHASE_COLORS["Origin"],
        PHASE_COLORS["Shipping"],
        PHASE_COLORS["Control"],
        PHASE_COLORS["Finance"],
        PHASE_COLORS["Delivery"],
        PHASE_COLORS["Closure"],
    ]
    left = 0.86
    for index, (label, color) in enumerate(zip(labels, colors)):
        width = 1.22
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + index * 1.42), Inches(6.35), Inches(width), Inches(0.38))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        if index < len(labels) - 1:
            arrow = slide.shapes.add_textbox(Inches(left + index * 1.42 + 1.24), Inches(6.355), Inches(0.18), Inches(0.28))
            p = arrow.text_frame.paragraphs[0]
            p.text = ">"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(101, 118, 135)


def add_cards(slide, cards, phase):
    count = len(cards)
    cols = 2 if count in (2, 4) else 3
    rows = 2 if count > cols else 1
    card_width = 5.55 if cols == 2 else 3.55
    card_height = 1.72 if rows == 1 else 1.52
    left_start = 0.92
    top_start = 1.82
    x_gap = 0.32
    y_gap = 0.28

    for index, (heading, body) in enumerate(cards):
        row = index // cols
        col = index % cols
        left = Inches(left_start + col * (card_width + x_gap))
        top = Inches(top_start + row * (card_height + y_gap))
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(card_width), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 252, 253)
        card.line.color.rgb = RGBColor(219, 227, 233)

        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.12), Inches(card_height))
        accent.fill.solid()
        accent.fill.fore_color.rgb = PHASE_COLORS[phase]
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.16), Inches(card_width - 0.35), Inches(0.34))
        p = title_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(25, 41, 58)

        body_box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.56), Inches(card_width - 0.35), Inches(card_height - 0.7))
        tf = body_box.text_frame
        tf.word_wrap = True
        lines = [line.strip() for line in body.splitlines() if line.strip()]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18 if len(lines) == 1 else 14)
            p.font.color.rgb = RGBColor(72, 90, 107)
            if len(lines) > 1:
                p.bullet = True
                p.space_after = Pt(6)


def add_diagram(slide, steps, phase, notes):
    phase_color = PHASE_COLORS[phase]
    start_left = 0.95
    top = 2.15
    step_width = 2.25 if len(steps) <= 4 else 1.7
    gap = 0.38

    for index, step in enumerate(steps):
        left = Inches(start_left + index * (step_width + gap))
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            Inches(top),
            Inches(step_width),
            Inches(1.05),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(249, 251, 253)
        box.line.color.rgb = phase_color
        box.line.width = Pt(2)
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = box.text_frame.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 46, 63)

        if index < len(steps) - 1:
            arrow = slide.shapes.add_textbox(
                Inches(start_left + index * (step_width + gap) + step_width + 0.02),
                Inches(top + 0.27),
                Inches(gap - 0.04),
                Inches(0.4),
            )
            p = arrow.text_frame.paragraphs[0]
            p.text = "→"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = phase_color

    note_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(4.2),
        Inches(11.2),
        Inches(1.25),
    )
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = RGBColor(241, 246, 249)
    note_box.line.fill.background()

    tf = note_box.text_frame
    tf.clear()
    for i, line in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(67, 84, 102)
        p.alignment = PP_ALIGN.CENTER
        if len(notes) > 1:
            p.space_after = Pt(6)


def add_snapshot_slide(slide, items):
    left_positions = [0.95, 6.27]
    top_positions = [1.95, 4.0]
    box_width = 5.1
    box_height = 1.55

    for index, (heading, body) in enumerate(items):
        left = Inches(left_positions[index % 2])
        top = Inches(top_positions[index // 2])
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            Inches(box_width),
            Inches(box_height),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 252, 253)
        card.line.color.rgb = RGBColor(214, 222, 229)

        stripe = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left,
            top,
            Inches(0.12),
            Inches(box_height),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(20, 84, 120)
        stripe.line.fill.background()

        add_text(slide, left + Inches(0.25), top + Inches(0.18), Inches(4.45), Inches(0.3), heading, 16, RGBColor(28, 44, 61), bold=True)
        add_text(slide, left + Inches(0.25), top + Inches(0.56), Inches(4.45), Inches(0.72), body, 13.5, RGBColor(82, 97, 113))

    ribbon = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.0),
        Inches(6.02),
        Inches(11.05),
        Inches(0.46),
    )
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = RGBColor(28, 49, 70)
    ribbon.line.fill.background()
    ribbon.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = ribbon.text_frame.paragraphs[0]
    p.text = "The platform acts as one control system across commercial setup, shipping execution, compliance, finance, inland delivery, and closure."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def add_control_rules_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, slide_data["phase"])
    add_title(slide, slide_data["title"], slide_data["subtitle"])
    add_phase_chip(slide, slide_data["phase"])

    center_x = 6.66
    center_y = 3.45
    core = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(center_x - 1.3),
        Inches(center_y - 0.48),
        Inches(2.6),
        Inches(0.96),
    )
    core.fill.solid()
    core.fill.fore_color.rgb = RGBColor(31, 49, 69)
    core.line.fill.background()
    core.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = core.text_frame.paragraphs[0]
    p.text = "SYSTEM\nENFORCEMENT"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    positions = [
        (1.05, 2.05),
        (8.05, 2.05),
        (1.05, 4.55),
        (8.05, 4.55),
    ]

    for (heading, body), (left, top) in zip(slide_data["cards"], positions):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(4.2),
            Inches(1.3),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(250, 252, 253)
        box.line.color.rgb = PHASE_COLORS[slide_data["phase"]]
        box.line.width = Pt(1.5)
        add_text(slide, Inches(left + 0.22), Inches(top + 0.18), Inches(3.75), Inches(0.28), heading, 15, RGBColor(27, 44, 61), bold=True)
        add_text(slide, Inches(left + 0.22), Inches(top + 0.55), Inches(3.75), Inches(0.42), body, 13.5, RGBColor(86, 100, 116))

    connectors = [
        ((5.25, 2.62), (5.85, 3.08)),
        ((8.05, 3.08), (7.95, 3.08)),
        ((5.25, 5.1), (5.85, 3.86)),
        ((8.05, 5.1), (7.95, 3.86)),
    ]
    for start, end in connectors:
        line = slide.shapes.add_connector(1, Inches(start[0]), Inches(start[1]), Inches(end[0]), Inches(end[1]))
        line.line.color.rgb = RGBColor(126, 141, 154)
        line.line.width = Pt(1.5)

    add_flow_strip(slide)


def add_closure_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, slide_data["phase"])
    add_title(slide, slide_data["title"], slide_data["subtitle"])
    add_phase_chip(slide, slide_data["phase"])

    ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4.7), Inches(1.9), Inches(3.9), Inches(3.9))
    ring.fill.solid()
    ring.fill.fore_color.rgb = RGBColor(245, 248, 250)
    ring.line.color.rgb = PHASE_COLORS[slide_data["phase"]]
    ring.line.width = Pt(4)

    inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(5.35), Inches(2.55), Inches(2.6), Inches(2.6))
    inner.fill.solid()
    inner.fill.fore_color.rgb = PHASE_COLORS[slide_data["phase"]]
    inner.line.fill.background()
    inner.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = inner.text_frame.paragraphs[0]
    p.text = "SHIPMENT\nCLOSED"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    labels = [
        ("Commercial complete", 1.05, 2.45),
        ("Documents cleared", 8.9, 2.45),
        ("Delivery confirmed", 1.1, 4.2),
        ("Container returned", 8.7, 4.2),
    ]
    for text, left, top in labels:
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.8), Inches(0.62))
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor(250, 252, 253)
        chip.line.color.rgb = RGBColor(196, 207, 216)
        chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = chip.text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = RGBColor(46, 63, 80)

    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(6.02), Inches(11.05), Inches(0.46))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = RGBColor(31, 49, 69)
    ribbon.line.fill.background()
    ribbon.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = ribbon.text_frame.paragraphs[0]
    p.text = "Closure is not just status marking. It confirms commercial, documentary, financial, delivery, and equipment return completion."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    add_flow_strip(slide)


def add_master_timeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, "Shipping")
    add_title(slide, "Master Timeline", "The shipment operating chain from quote to final closure")
    add_phase_chip(slide, "Shipping")

    rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(2.58), Inches(11.42), Inches(0.32))
    rail.fill.solid()
    rail.fill.fore_color.rgb = RGBColor(220, 229, 236)
    rail.line.fill.background()

    start_left = 1.0
    top = 1.95
    width = 1.38
    gap = 0.28

    for index, (phase, lines) in enumerate(MASTER_TIMELINE):
        left = start_left + index * (width + gap)
        color = PHASE_COLORS[phase]

        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + 0.48), Inches(2.42), Inches(0.42), Inches(0.42))
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()

        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(1.18),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        card.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = card.text_frame.paragraphs[0]
        p.text = phase
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        detail = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left - 0.05),
            Inches(3.1),
            Inches(width + 0.1),
            Inches(1.52),
        )
        detail.fill.solid()
        detail.fill.fore_color.rgb = RGBColor(251, 252, 253)
        detail.line.color.rgb = color
        detail.line.width = Pt(1.5)

        tf = detail.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate([item.strip() for item in lines.splitlines() if item.strip()]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(38, 54, 71)
            p.space_after = Pt(6)

        if index < len(MASTER_TIMELINE) - 1:
            arrow = slide.shapes.add_textbox(
                Inches(left + width + 0.02),
                Inches(2.14),
                Inches(gap - 0.04),
                Inches(0.5),
            )
            p = arrow.text_frame.paragraphs[0]
            p.text = "→"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = RGBColor(110, 126, 141)

    summary = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(5.05),
        Inches(11.15),
        Inches(0.95),
    )
    summary.fill.solid()
    summary.fill.fore_color.rgb = RGBColor(31, 49, 69)
    summary.line.fill.background()
    summary.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = summary.text_frame.paragraphs[0]
    p.text = "A shipment moves through seven controlled phases: commercial setup, origin readiness, shipping execution, release control, financial settlement, inland delivery, and formal closure."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    add_flow_strip(slide)


def add_slide(prs, slide_data, title=False):
    if slide_data.get("snapshot"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_full_background(slide, slide_data["phase"])
        add_title(slide, slide_data["title"], slide_data["subtitle"])
        add_phase_chip(slide, slide_data["phase"])
        add_snapshot_slide(slide, slide_data["snapshot"])
        add_flow_strip(slide)
        return
    if slide_data["title"] == "Control Rules":
        add_control_rules_slide(prs, slide_data)
        return
    if slide_data["title"] == "Closure":
        add_closure_slide(prs, slide_data)
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, slide_data["phase"])
    add_title(slide, slide_data["title"], slide_data["subtitle"])
    add_phase_chip(slide, slide_data["phase"])
    if "diagram" in slide_data:
        add_diagram(slide, slide_data["diagram"], slide_data["phase"], slide_data.get("notes", []))
    else:
        add_cards(slide, slide_data["cards"], slide_data["phase"])
    add_flow_strip(slide)

    if title:
        hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(5.0), Inches(11.4), Inches(0.82))
        hero.fill.solid()
        hero.fill.fore_color.rgb = RGBColor(233, 241, 246)
        hero.line.fill.background()
        p = hero.text_frame.paragraphs[0]
        p.text = "One controlled operating chain across commercial, document, bank, port, finance, dispatch, and return"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(24, 52, 78)


def add_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, "Shipping")
    add_title(slide, "Detailed Process Flow", "31 operational steps from quote to shipment closure")
    add_phase_chip(slide, "Shipping")

    start_left = 0.86
    top = 1.82
    card_width = 1.58
    gap = 0.14

    for index, (phase, steps) in enumerate(FLOW_PHASES):
        left = Inches(start_left + index * (card_width + gap))
        color = PHASE_COLORS[phase]

        band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            Inches(top),
            Inches(card_width),
            Inches(0.46),
        )
        band.fill.solid()
        band.fill.fore_color.rgb = color
        band.line.fill.background()
        band.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = band.text_frame.paragraphs[0]
        p.text = phase
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        lane = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            Inches(top + 0.58),
            Inches(card_width),
            Inches(4.8),
        )
        lane.fill.solid()
        lane.fill.fore_color.rgb = RGBColor(251, 252, 253)
        lane.line.color.rgb = RGBColor(217, 225, 232)

        cursor_y = top + 0.76
        step_height = 0.52
        inner_width = card_width - 0.22

        for step_number, step in enumerate(steps, start=sum(len(items) for _, items in FLOW_PHASES[:index]) + 1):
            step_box = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(start_left + index * (card_width + gap) + 0.11),
                Inches(cursor_y),
                Inches(inner_width),
                Inches(step_height),
            )
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = RGBColor(246, 248, 250)
            step_box.line.color.rgb = color
            step_box.line.width = Pt(1.2)
            step_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            tf = step_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"{step_number}. {step}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10.5)
            p.font.bold = True
            p.font.color.rgb = RGBColor(39, 56, 73)

            cursor_y += step_height + 0.1

            if step_number < sum(len(items) for _, items in FLOW_PHASES[: index + 1]):
                arrow = slide.shapes.add_textbox(
                    Inches(start_left + index * (card_width + gap) + 0.64),
                    Inches(cursor_y - 0.03),
                    Inches(0.3),
                    Inches(0.18),
                )
                p = arrow.text_frame.paragraphs[0]
                p.text = "↓"
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = color

        if index < len(FLOW_PHASES) - 1:
            connector = slide.shapes.add_textbox(
                Inches(start_left + index * (card_width + gap) + card_width - 0.01),
                Inches(top + 2.52),
                Inches(gap),
                Inches(0.36),
            )
            p = connector.text_frame.paragraphs[0]
            p.text = "→"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(109, 124, 138)

    summary = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(6.0),
        Inches(11.15),
        Inches(0.48),
    )
    summary.fill.solid()
    summary.fill.fore_color.rgb = RGBColor(239, 244, 247)
    summary.line.fill.background()
    summary.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = summary.text_frame.paragraphs[0]
    p.text = "Every shipment moves across seven controlled lanes, with release and closure blocked until the required upstream steps are complete."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(34, 52, 71)

    add_flow_strip(slide)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, SLIDES[0])
    add_master_timeline_slide(prs)
    for slide_data in SLIDES[1:]:
        add_slide(prs, slide_data)
    add_flow_slide(prs)

    prs.save(OUTPUT_FILE)
    print(f"PPT generated successfully: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
